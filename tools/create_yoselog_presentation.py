from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation"
PPTX_PATH = OUT / "YoseLog_Final_Presentation.pptx"
SCRIPT_PATH = OUT / "YoseLog_10min_Speaker_Script.md"

WIDE = (13.333, 7.5)

COLORS = {
    "ink": RGBColor(24, 33, 30),
    "muted": RGBColor(93, 111, 104),
    "mint": RGBColor(47, 201, 142),
    "mint_dark": RGBColor(20, 122, 88),
    "canvas": RGBColor(237, 243, 240),
    "panel": RGBColor(255, 255, 255),
    "line": RGBColor(206, 220, 214),
    "gold": RGBColor(224, 168, 59),
    "rose": RGBColor(205, 96, 111),
    "charcoal": RGBColor(21, 29, 32),
}


def set_text(shape, text, size=20, color="ink", bold=False, align=None):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = COLORS[color]
    p.font.bold = bold
    p.font.name = "Aptos"
    if align:
        p.alignment = align


def add_text(slide, text, x, y, w, h, size=20, color="ink", bold=False, align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size=size, color=color, bold=bold, align=align)
    return shape


def add_rect(slide, x, y, w, h, fill="panel", line="line", radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS[line]
    shape.line.width = Pt(1)
    return shape


def add_header(slide, title, subtitle=None):
    add_text(slide, title, 0.65, 0.36, 8.6, 0.45, size=22, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.66, 0.83, 8.6, 0.3, size=9.5, color="muted")
    add_text(slide, "YoseLog", 11.45, 0.42, 1.1, 0.25, size=10, color="mint_dark", bold=True, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.18), Inches(12.0), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()


def add_bullets(slide, items, x, y, w, h, size=16, color="ink"):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Aptos"
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(7)
    return shape


def add_label(slide, text, x, y, w, color="mint_dark"):
    pill = add_rect(slide, x, y, w, 0.34, fill="canvas", line="line")
    set_text(pill, text, size=9, color=color, bold=True, align=PP_ALIGN.CENTER)
    pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return pill


def add_arrow(slide, x1, y1, x2, y2, color="mint_dark"):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COLORS[color]
    line.line.width = Pt(2)
    line.line.end_arrowhead = True
    return line


def add_mlp_diagram(slide, x, y):
    columns = [
        ("Input\n66 features", x, y + 0.45, 4, "canvas"),
        ("Hidden\n256", x + 2.1, y + 0.15, 5, "gold"),
        ("Hidden\n128", x + 4.2, y + 0.35, 4, "rose"),
        ("Output\n40 poses", x + 6.25, y + 0.85, 3, "mint"),
    ]

    node_positions = []
    for label, cx, cy, count, color in columns:
        positions = []
        spacing = 0.58
        start = cy
        add_text(slide, label, cx - 0.55, y - 0.05, 1.4, 0.52, size=9.5, bold=True, align=PP_ALIGN.CENTER)
        for i in range(count):
            yy = start + i * spacing
            node = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(cx),
                Inches(yy),
                Inches(0.42),
                Inches(0.42),
            )
            node.fill.solid()
            node.fill.fore_color.rgb = COLORS[color]
            node.line.color.rgb = COLORS["mint_dark"]
            node.line.width = Pt(1)
            positions.append((cx + 0.21, yy + 0.21))
        node_positions.append(positions)

    for left_nodes, right_nodes in zip(node_positions, node_positions[1:]):
        for x1, y1 in left_nodes:
            for x2, y2 in right_nodes:
                conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
                conn.line.color.rgb = COLORS["line"]
                conn.line.width = Pt(0.6)

    add_text(slide, "Each neuron combines weighted landmark features and learns nonlinear pose patterns.", x + 0.1, y + 3.55, 6.8, 0.35, size=10.5, color="muted", align=PP_ALIGN.CENTER)


def metric_card(slide, label, value, note, x, y, w=2.45):
    card = add_rect(slide, x, y, w, 1.2)
    add_text(slide, label, x + 0.18, y + 0.16, w - 0.36, 0.2, size=8.5, color="muted", bold=True)
    add_text(slide, value, x + 0.18, y + 0.43, w - 0.36, 0.32, size=20, color="mint_dark", bold=True)
    add_text(slide, note, x + 0.18, y + 0.85, w - 0.36, 0.22, size=8.5, color="muted")
    return card


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(WIDE[0])
    prs.slide_height = Inches(WIDE[1])
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    bg = add_rect(s, 0, 0, 13.333, 7.5, fill="charcoal", line="charcoal", radius=False)
    add_text(s, "YoseLog", 0.72, 0.65, 4.2, 0.8, size=42, color="panel", bold=True)
    add_text(s, "Real-time yoga pose classification, balance scoring, and automatic session logging", 0.78, 1.55, 8.8, 0.55, size=17, color="canvas")
    add_label(s, "10-minute project defense", 0.8, 2.3, 2.25, color="mint_dark")
    for x, y, text in [
        (0.9, 4.75, "MediaPipe Pose"),
        (3.25, 4.75, "ML classifier"),
        (5.6, 4.75, "Stable hold detection"),
        (8.25, 4.75, "Web dashboard"),
    ]:
        add_rect(s, x, y, 1.9, 0.78, fill="panel", line="line")
        add_text(s, text, x + 0.15, y + 0.25, 1.6, 0.2, size=10, color="ink", bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "Presenter: Dahyun Eom", 0.82, 6.65, 3.2, 0.25, size=11, color="canvas")

    # 2
    s = prs.slides.add_slide(blank)
    add_header(s, "Problem: most yoga apps assume the pose first", "My goal was free-form session logging: let the user move naturally while the app detects and logs poses automatically.")
    add_rect(s, 0.75, 1.75, 3.45, 3.65, fill="canvas")
    add_text(s, "Typical flow", 1.0, 2.05, 2.8, 0.3, size=16, bold=True)
    add_bullets(s, ["User selects a target pose", "App checks if the pose matches", "Scoring/feedback is tied to that pre-selected pose", "Hard-coded rules often drive pose-specific feedback"], 1.0, 2.6, 2.65, 2.3, size=12.4)
    add_rect(s, 4.95, 1.75, 3.45, 3.65, fill="panel")
    add_text(s, "My observation", 5.2, 2.05, 2.8, 0.3, size=16, bold=True)
    add_bullets(s, ["A real yoga session is not always pre-planned", "The user may switch poses freely", "Manual pose selection interrupts the session", "Automatic logging would be more natural"], 5.2, 2.6, 2.65, 2.3, size=12.4)
    add_rect(s, 9.15, 1.75, 3.45, 3.65, fill="canvas")
    add_text(s, "Final product", 9.4, 2.05, 2.8, 0.3, size=16, bold=True)
    add_bullets(s, ["Confirmed classifier result", "Automatic session table", "Pose durations + totals", "Exportable CSV for review"], 9.4, 2.6, 2.65, 2.3, size=12.4)

    # 3
    s = prs.slides.add_slide(blank)
    add_header(s, "System architecture", "The frontend does not invent labels; it displays the same confirmed result from the live classifier loop.")
    stages = [
        ("Camera", "OpenCV VideoCapture"),
        ("Pose landmarks", "MediaPipe Pose\n33 body points"),
        ("Classifier", "66 x/y features\nscikit-learn model"),
        ("Session engine", "stability + hold time\nbalance score"),
        ("Website", "Flask API + MJPEG\nbrowser dashboard"),
    ]
    x = 0.75
    for i, (title, body) in enumerate(stages):
        add_rect(s, x, 2.1, 2.0, 1.45, fill="panel")
        add_text(s, title, x + 0.15, 2.32, 1.7, 0.25, size=14, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, body, x + 0.15, 2.74, 1.7, 0.45, size=9.6, color="muted", align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            add_arrow(s, x + 2.05, 2.82, x + 2.55, 2.82)
        x += 2.55
    add_text(s, "Design decision", 0.85, 4.65, 1.6, 0.25, size=12, color="mint_dark", bold=True)
    add_bullets(s, ["The app was aligned back to classify_live.py after testing showed frontend-only gates changed pose results.", "This made the web UI a presentation layer over the classifier, not a second classifier."], 0.85, 5.0, 11.5, 0.9, size=13, color="ink")

    # 4
    s = prs.slides.add_slide(blank)
    add_header(s, "Why I chose an MLP", "The input is already structured landmark data, so a small neural network is a practical classifier.")
    add_rect(s, 0.75, 1.48, 7.55, 4.95, fill="panel")
    add_mlp_diagram(s, 1.05, 2.02)
    add_rect(s, 8.75, 1.48, 3.8, 4.95, fill="canvas")
    add_text(s, "Why MLP fits", 9.05, 1.82, 2.6, 0.3, size=16, bold=True)
    add_bullets(s, ["Input is not raw pixels", "MediaPipe already gives 66 numeric x/y features", "MLP learns nonlinear joint relationships", "Lighter than training a CNN", "Easier than hand-coding angle rules for 40 poses"], 9.05, 2.35, 3.05, 2.55, size=11.7)
    add_text(s, "Defense point: MediaPipe handles perception; the MLP handles classification from clean landmark features.", 0.95, 6.65, 11.2, 0.3, size=13.5, color="mint_dark", bold=True, align=PP_ALIGN.CENTER)

    # 5
    s = prs.slides.add_slide(blank)
    add_header(s, "Classifier pipeline", "I trained the pose classifier from downloaded Yoga-82 images using MediaPipe keypoints and an MLP model.")
    metric_card(s, "POSE CLASSES", "40", "downloaded pose folders", 0.75, 1.55)
    metric_card(s, "SAMPLES", "11,566", "keypoint rows extracted", 3.55, 1.55)
    metric_card(s, "TRAIN / TEST", "9,252 / 2,314", "stratified split", 6.35, 1.55)
    metric_card(s, "ACCURACY", "86.9%", "held-out test set", 9.15, 1.55)
    add_rect(s, 0.8, 3.25, 11.75, 1.35, fill="canvas")
    add_text(s, "Training process", 1.05, 3.5, 2.0, 0.25, size=13, bold=True)
    add_text(s, "download images → original + flipped keypoint extraction → 66-feature keypoints.csv → train_model_mlp.py → pose_model.pkl", 1.05, 3.95, 10.95, 0.32, size=14.4, color="ink")
    add_rect(s, 0.8, 5.0, 5.55, 1.05, fill="panel")
    add_text(s, "Accuracy key", 1.05, 5.2, 1.5, 0.2, size=11, color="muted", bold=True)
    add_text(s, "Horizontal flip augmentation doubled orientation coverage for left/right body views", 1.05, 5.55, 4.8, 0.25, size=11.8)
    add_rect(s, 7.0, 5.0, 5.55, 1.05, fill="panel")
    add_text(s, "Model", 7.25, 5.2, 1.2, 0.2, size=11, color="muted", bold=True)
    add_text(s, "MLPClassifier hidden layers (256, 128), trained on original + flipped landmarks", 7.25, 5.55, 4.8, 0.25, size=11.8)

    # 6
    s = prs.slides.add_slide(blank)
    add_header(s, "Model tuning experiment", "Changing hidden-layer size changes accuracy, but most reasonable MLPs clustered near 86-87%.")
    graph_path = ROOT / "analysis" / "mlp_architecture_accuracy.png"
    if graph_path.exists():
        s.shapes.add_picture(str(graph_path), Inches(0.7), Inches(1.45), width=Inches(7.15))
    add_rect(s, 8.25, 1.48, 4.15, 4.85, fill="canvas")
    add_text(s, "What the graph shows", 8.55, 1.82, 3.2, 0.3, size=16, bold=True)
    add_bullets(s, [
        "Best single result: (128,) at 87.0%",
        "(256,128) reached 86.9%, nearly tied with the best",
        "(256,128) had slightly stronger macro-F1 on the imbalanced 40-pose dataset",
        "Very deep/wide is not automatically better",
    ], 8.55, 2.35, 3.35, 2.25, size=11.7)
    add_text(s, "Defense point: I kept (256,128) because it balances accuracy, capacity, and per-class performance.", 8.55, 5.35, 3.35, 0.55, size=11.2, color="muted")

    # 7
    s = prs.slides.add_slide(blank)
    add_header(s, "Session logging converts predictions into usable events", "A pose is shown/logged only after it survives the stability rule.")
    add_rect(s, 0.9, 1.65, 3.2, 3.9, fill="panel")
    add_text(s, "Stability rule", 1.15, 1.95, 2.6, 0.3, size=15, bold=True)
    add_bullets(s, ["Compare current smoothed keypoints to previous frame", "Movement = mean absolute keypoint difference", "If movement < 0.05 and label is unchanged, count stable frame", "30 stable frames ≈ 1 second"], 1.15, 2.45, 2.55, 2.4, size=11.8)
    add_rect(s, 5.0, 1.65, 3.2, 3.9, fill="canvas")
    add_text(s, "Logging rule", 5.25, 1.95, 2.6, 0.3, size=15, bold=True)
    add_bullets(s, ["Keep accumulating hold time for the same pose", "When pose changes, log the previous pose", "Only log if duration ≥ 1 second", "Show recent entries in the UI"], 5.25, 2.45, 2.55, 2.4, size=11.8)
    add_rect(s, 9.1, 1.65, 3.2, 3.9, fill="panel")
    add_text(s, "Why it matters", 9.35, 1.95, 2.6, 0.3, size=15, bold=True)
    add_bullets(s, ["Reduces flicker", "Prevents accidental one-frame logs", "Turns recognition into a session summary", "Makes export meaningful"], 9.35, 2.45, 2.55, 2.4, size=11.8)

    # 8
    s = prs.slides.add_slide(blank)
    add_header(s, "Balance score estimates steadiness, not pose correctness", "Balance is computed from body sway during a held pose.")
    add_rect(s, 0.85, 1.6, 5.4, 4.65, fill="canvas")
    add_text(s, "What is measured", 1.15, 1.95, 3.0, 0.3, size=16, bold=True)
    add_bullets(s, ["Torso center movement over the hold", "Shoulder tilt variation over the hold", "Normalized by torso length so camera distance has less effect"], 1.15, 2.55, 4.55, 1.35, size=13)
    add_text(s, "score = 100 - sway penalty - tilt penalty", 1.15, 4.45, 4.55, 0.32, size=17, color="mint_dark", bold=True)
    add_text(s, "Higher score = steadier hold", 1.15, 5.05, 4.55, 0.3, size=12, color="muted")
    add_rect(s, 7.05, 1.6, 5.4, 4.65, fill="panel")
    add_text(s, "Defense distinction", 7.35, 1.95, 3.0, 0.3, size=16, bold=True)
    add_bullets(s, ["It does not grade yoga form accuracy", "It does not compare against an ideal pose template", "It only measures how much the detected body landmarks wobble during a stable hold"], 7.35, 2.55, 4.55, 1.6, size=13)
    add_text(s, "This makes the feature honest: balance is a stability metric, not a coach.", 7.35, 5.05, 4.6, 0.45, size=13, color="muted")

    # 9
    s = prs.slides.add_slide(blank)
    add_header(s, "Frontend integration", "The website makes the classifier usable during a yoga session.")
    add_rect(s, 0.75, 1.55, 7.0, 4.75, fill="charcoal", line="charcoal")
    add_text(s, "Live session dashboard", 1.1, 1.9, 3.2, 0.35, size=18, color="panel", bold=True)
    for i, text in enumerate(["Live camera stream", "Confirmed classifier result", "Hold timer + progress", "Balance score", "Session log + CSV export"]):
        add_rect(s, 1.1, 2.55 + i * 0.58, 5.7, 0.38, fill="panel", line="line")
        add_text(s, text, 1.28, 2.64 + i * 0.58, 5.2, 0.15, size=10.5, color="ink")
    add_rect(s, 8.25, 1.55, 4.35, 4.75, fill="canvas")
    add_text(s, "Key integration choice", 8.55, 1.9, 3.4, 0.35, size=16, bold=True)
    add_bullets(s, ["The browser receives the same confirmed result as classify_live.py", "The backend owns camera + model inference", "The frontend polls session state and renders the stream", "Port moved to 5001 to avoid macOS AirPlay conflict"], 8.55, 2.45, 3.45, 2.5, size=12.2)

    # 10
    s = prs.slides.add_slide(blank)
    add_header(s, "My contribution", "I moved the project from a classifier script to a usable session-logging product.")
    rows = [
        ("Training", "Downloaded data for 40 yoga poses and extracted original + horizontally flipped keypoints"),
        ("Model", "Wrote train_model_mlp.py and trained an MLP classifier with 86.9% test accuracy"),
        ("Product", "Built Flask website, live stream, state API, dashboard, and CSV export"),
        ("Reliability", "Aligned frontend result to classify_live.py after debugging mismatches"),
        ("Feature", "Added balance scoring from held-pose landmark sway"),
    ]
    y = 1.55
    for tag, desc in rows:
        add_rect(s, 0.9, y, 2.0, 0.55, fill="canvas")
        add_text(s, tag, 1.1, y + 0.18, 1.55, 0.18, size=10, color="mint_dark", bold=True, align=PP_ALIGN.CENTER)
        add_rect(s, 3.15, y, 8.9, 0.55, fill="panel")
        add_text(s, desc, 3.35, y + 0.17, 8.45, 0.18, size=12.5)
        y += 0.75
    add_text(s, "Contribution claim: I contributed the dataset/model training path and the product layer that logs free-form sessions.", 0.95, 6.05, 11.2, 0.35, size=15, color="mint_dark", bold=True)

    # 11
    s = prs.slides.add_slide(blank)
    add_header(s, "Limitations and defense answers", "The system is honest about what it can and cannot claim.")
    add_rect(s, 0.8, 1.45, 5.75, 4.9, fill="panel")
    add_text(s, "Known limitations", 1.1, 1.8, 2.6, 0.3, size=16, bold=True)
    add_bullets(s, ["Closed-set model can still confuse visually similar poses", "No explicit no-pose/resting training class", "Balance score measures steadiness, not correctness", "Camera angle and body visibility affect landmarks"], 1.1, 2.35, 4.7, 2.2, size=12.5)
    add_rect(s, 6.95, 1.45, 5.75, 4.9, fill="canvas")
    add_text(s, "How I defend it", 7.25, 1.8, 2.6, 0.3, size=16, bold=True)
    add_bullets(s, ["Separated raw classifier testing from web logging", "Used stable-frame confirmation to reduce flicker", "Kept frontend aligned with the trusted live classifier", "Designed next steps: add no-pose data and pose-correctness scoring"], 7.25, 2.35, 4.7, 2.2, size=12.5)

    # 12
    s = prs.slides.add_slide(blank)
    add_header(s, "Final takeaway", "YoseLog turns pose detection into a session record a user can actually review.")
    add_text(s, "What I built", 0.95, 1.6, 2.1, 0.3, size=15, color="mint_dark", bold=True)
    add_bullets(s, ["Real-time yoga pose classification", "Automatic duration logging", "Balance scoring during stable holds", "Browser dashboard and export"], 0.95, 2.1, 4.2, 2.0, size=14)
    add_text(s, "Why it matters", 6.3, 1.6, 2.3, 0.3, size=15, color="mint_dark", bold=True)
    add_bullets(s, ["It reduces manual tracking", "It makes model output understandable", "It creates a foundation for future coaching features"], 6.3, 2.1, 4.9, 1.55, size=14)
    add_rect(s, 1.0, 5.45, 11.1, 0.85, fill="charcoal", line="charcoal")
    add_text(s, "Defense thesis: my main contribution was building the reliable session layer around the classifier.", 1.35, 5.75, 10.3, 0.25, size=17, color="panel", bold=True, align=PP_ALIGN.CENTER)

    prs.save(PPTX_PATH)


SCRIPT = """# YoseLog 10-Minute Presentation Script

## Slide 1 — Title
Hi, my project is YoseLog, a real-time yoga session logger. The goal is not only to recognize poses from a webcam, but to turn those predictions into a usable workout/session record. The final system classifies yoga poses, confirms stable holds, estimates balance, logs duration automatically, and shows everything in a web dashboard.

## Slide 2 — Problem
Many yoga applications are built around a pre-selected pose. The user chooses the pose first, then the app measures how well the user matches that pose and gives scoring or feedback. That can work for pose correction, but it is not ideal for a free yoga session, because the user has to keep telling the app what pose they are about to do. My idea was: what if the user can just do the session naturally, and the application automatically detects the pose, confirms it, and logs the whole session?

## Slide 3 — System Architecture
The system starts with OpenCV capturing webcam frames. MediaPipe Pose extracts 33 body landmarks, which become 66 x/y features. Those features go into the trained scikit-learn model. Then the session engine applies the same stability logic from classify_live.py, computes hold duration and balance, and sends the state to the Flask frontend. A key design decision was making the frontend display the confirmed result from the classifier loop, instead of inventing a separate frontend interpretation.

## Slide 4 — Why I Chose an MLP
I chose an MLP because my model input is not raw images. MediaPipe already converts each image into structured body landmarks, so the classifier receives 66 numeric features: x and y for 33 points. For this type of tabular landmark data, an MLP is a practical choice because it can learn nonlinear relationships between joints, such as how shoulders, hips, knees, and ankles are positioned relative to each other. It is also much lighter than training a CNN on raw images, and it is easier to iterate on for this project. Compared with hard-coded angle rules, the MLP can learn from examples and is easier to expand to 40 poses.

## Slide 5 — Classifier Pipeline
For the classifier, I downloaded images for 40 yoga poses and used MediaPipe Pose to extract 33 body landmarks from each image. Each sample becomes 66 numeric features, x and y for each landmark, and those rows are saved in keypoints.csv. A key improvement was that during keypoint extraction, I also created a horizontally flipped version of each image using cv2.flip and extracted landmarks from that flipped image too. This gave the model more left/right orientation coverage and helped bring the accuracy to around 87 percent. I wrote train_model_mlp.py to train an MLPClassifier with hidden layers of 256 and 128. The final dataset had 11,566 samples, split into 9,252 training samples and 2,314 testing samples. The model reached 86.9 percent accuracy on the held-out test set, and the trained model is saved as pose_model.pkl. At runtime, webcam landmarks are converted into the same 66-feature format, smoothed over five frames, and passed into model.predict.

## Slide 6 — Model Tuning Experiment
I also tested different MLP hidden-layer architectures. This helped answer the question of whether 256 and 128 was arbitrary. The results showed that most reasonable models clustered around 86 to 87 percent accuracy. The single hidden layer with 128 neurons was slightly highest at about 87.0 percent, while 256 and 128 reached about 86.9 percent, almost tied. I kept 256 and 128 because the accuracy difference was tiny, and 256 and 128 had slightly stronger macro-F1. That matters because my 40-pose dataset is imbalanced, so macro-F1 better reflects performance across rare classes as well as common classes. Larger models like 512 and 256 did not clearly improve accuracy, and the three-layer model performed worse. So my conclusion is that hidden-layer size matters, but more layers or more neurons are not automatically better.

## Slide 7 — Session Logging Logic
The session logger makes predictions usable. It compares the current smoothed keypoints to the previous smoothed keypoints. If the predicted pose is the same and the movement is under 0.05, the stable frame counter increases. The system requires 30 stable frames, which is about one second. Only then is the pose treated as a confirmed hold. When the pose changes, the previous pose is logged if it lasted at least one second.

## Slide 8 — Balance Scoring
Balance scoring is separate from classification. It does not judge whether the yoga form is correct. It estimates steadiness during a held pose. I store the stable keypoints during the hold, then measure torso center movement and shoulder tilt variation. Less sway means a higher balance score. This is important to explain clearly: balance score is a stability metric, not a yoga-correction metric.

## Slide 9 — Frontend Integration
The frontend turns the classifier into an application. It shows the live video stream, confirmed classifier result, hold timer, stability progress, balance score, and session log. The backend owns camera and inference; the browser only displays the state. During testing, I found that extra frontend rules made results differ from classify_live.py, so I aligned the web result back to the same classifier mechanism.

## Slide 10 — My Contribution
My contribution has two parts: the classifier work and the product work. For the classifier, I downloaded the dataset for 40 poses, extracted MediaPipe keypoints, added horizontal flip augmentation during extraction, wrote train_model_mlp.py, trained the MLP model, and evaluated it at 86.9 percent accuracy. For the product, I built the Flask web app, state API, live stream, dashboard, CSV export, and balance scoring. I also debugged the difference between raw model behavior and frontend behavior, and restored the frontend to follow the reliable classify_live.py path.

## Slide 11 — Limitations and Defense
There are still limitations. Similar poses can be confused, especially because this is a closed-set classifier. There is no explicit no-pose/resting class yet. Balance measures stillness, not correctness. Camera angle and body visibility matter. My defense is that I separated raw classifier testing from web logging, used stable-frame confirmation, kept the frontend aligned with the classifier, and identified clear next steps: add no-pose training data, collect more examples for confusing poses, and add pose-correctness scoring.

## Slide 12 — Final Takeaway
The final takeaway is that YoseLog turns pose detection into a session record. The user can see the confirmed pose, duration, balance score, and session log automatically. My main contribution is the reliable session layer around the classifier: taking a model output and building the application logic needed for a usable yoga logging experience.

# Professor Defense Q&A

## Q: Why does the classifier sometimes choose a pose when the user is resting?
Because the model is closed-set. It only knows the trained pose classes, so it must pick the nearest one. A better future version would include a no-pose/resting class.

## Q: Is the balance score measuring pose correctness?
No. It measures steadiness during a held pose using torso center sway and shoulder tilt variation. Correctness would require comparing joint angles to an ideal pose template.

## Q: Why did you align the frontend to classify_live.py?
Because the frontend should not become a second classifier. When extra UI gates changed results, I made the backend follow the same mechanism as the trusted live classifier loop.

## Q: What did you personally contribute?
I downloaded and prepared the 40-pose dataset, extracted MediaPipe keypoints, added flipped-version augmentation, wrote the MLP training script, trained/evaluated the classifier, and then built the web application layer: session state API, live dashboard, CSV export, balance scoring, raw classifier tester, and the debugging/alignment work that made the frontend reflect the classifier correctly.

## Q: What would you improve next?
I would add a no-pose/resting dataset class, collect more examples for visually similar poses, add per-pose correctness scoring using joint angles, and evaluate accuracy with live webcam test videos instead of only image-based validation.

## Q: Why use hidden layers of 256 and 128 if 128 was slightly higher in the experiment?
The difference was very small: around 87.0 percent versus 86.9 percent. I kept 256 and 128 because it had almost the same accuracy but slightly better macro-F1. Since the dataset is imbalanced, macro-F1 is important because it gives rare poses equal weight. I would explain it as a balanced choice, not the only possible choice. Future work could select the final architecture by cross-validation.

## Q: Why did you choose MLP instead of a CNN?
Because the input to my classifier is not raw images. MediaPipe already extracts structured landmark features, so the model receives 66 numeric values. An MLP is appropriate for this compact tabular feature vector, while a CNN would be more useful if I were training directly on raw pixels.
"""


def write_script():
    SCRIPT_PATH.write_text(SCRIPT, encoding="utf-8")


if __name__ == "__main__":
    OUT.mkdir(exist_ok=True)
    build_deck()
    write_script()
    print(PPTX_PATH)
    print(SCRIPT_PATH)
